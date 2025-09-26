from pathlib import Path
import sys

from pptx import Presentation

p = sys.argv[1]
shift_offset = int(sys.argv[2]) # Usually -7 or +7

# WARNING: contains non-printable characters!
SHIFT_TABLE = """
%^&1234567!@#$
TYUqwertyuQWER
GHJasdfghjASDF
BNMzxcvbnmZXCV





""".replace('\n', '')

def shift_note(text, shift):
    result = []
    for char in text:
        if char in SHIFT_TABLE:
            index = SHIFT_TABLE.index(char)
            new_index = index + shift
            # each line has 35 characters
            # check if index and new_index is in the same line
            if (index // 35) == (new_index // 35) and 0 <= new_index < len(SHIFT_TABLE):
                result.append(SHIFT_TABLE[new_index])
            else:
                raise ValueError(f"Shift out of bounds for character: {char}, line: {index // 35 + 1}, col: {index % 35 + 1}")
        else:
            result.append(char)  # Not in table, keep original
    return ''.join(result)


prs = Presentation(p)
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                # text_runs.append((run.font.name, run.text))
                if run.font.name == "SimpMusic Base":
                    run.text = shift_note(run.text, shift_offset)

prs.save(Path(p).with_stem(Path(p).stem + "_shifted"))
